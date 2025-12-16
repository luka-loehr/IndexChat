#!/usr/bin/env python3
"""
IndexChat Cloud-First Indexer
Indexes documents using OpenAI and Hugging Face APIs.
"""

import argparse
import os
import struct
import sqlite3
import time
import base64
from pathlib import Path
import mimetypes
import requests
import cv2
import numpy as np

import tiktoken
from openai import OpenAI
from dotenv import load_dotenv
from PIL import Image

# Optional imports for document handling
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# Load environment variables
root_env_path = Path(__file__).parent.parent / ".env"

if root_env_path.exists():
    load_dotenv(dotenv_path=root_env_path, override=True)
else:
    load_dotenv(override=True)

# Configuration
INPUT_DIR = Path(__file__).parent.parent / "input"
DB_PATH = Path(__file__).parent / "database.sqlite"

# Models
OPENAI_TEXT_EMBED_MODEL = "text-embedding-3-large"
HF_IMAGE_MODEL = "openai/clip-vit-base-patch32"
HF_AUDIO_MODEL = "laion/clap-htsat-unfused" 

# Dimensions
TEXT_EMBED_DIM = 3072
IMAGE_EMBED_DIM = 512
AUDIO_EMBED_DIM = 512
                      
CHUNK_SIZE = 800
CHUNK_OVERLAP = 100

# File Extensions
DOC_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md"}
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a"}
VIDEO_EXTENSIONS = {".mp4", ".mov", ".avi"}

def get_openai_client() -> OpenAI:
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise ValueError("OPENAI_API_KEY environment variable is required.")
    return OpenAI(api_key=api_key)

def get_hf_headers():
    token = os.getenv("HUGGINGFACE_API_KEY")
    if not token:
        print("⚠️ HUGGINGFACE_API_KEY not found. Image/Audio embeddings will fail.")
        return {}
    return {"Authorization": f"Bearer {token}"}

def query_hf_api(model_id, data, retries=3):
    api_url = f"https://api-inference.huggingface.co/models/{model_id}"
    headers = get_hf_headers()
    
    for i in range(retries):
        try:
            response = requests.post(api_url, headers=headers, json=data)
            if response.status_code == 200:
                return response.json()
            else:
                err = response.json()
                if "error" in err and "loading" in err["error"].lower():
                    wait_time = err.get("estimated_time", 20)
                    print(f"Model {model_id} loading, waiting {wait_time}s...")
                    time.sleep(wait_time)
                    continue
                print(f"Error querying HF API ({model_id}): {response.status_code} - {response.text}")
        except Exception as e:
            print(f"Exception querying HF API: {e}")
        time.sleep(2)
    return None

def get_hf_image_embedding(image_path=None, image_data=None):
    # Support both path and direct bytes
    headers = get_hf_headers()
    api_url = f"https://api-inference.huggingface.co/pipeline/feature-extraction/{HF_IMAGE_MODEL}"
    
    data = image_data
    if image_path:
        with open(image_path, "rb") as f:
            data = f.read()
            
    if not data: return None
        
    response = requests.post(api_url, headers=headers, data=data)
    if response.status_code != 200:
        print(f"Error getting image embedding: {response.text}")
        return None
        
    return response.json()

def get_hf_audio_embedding(audio_path: Path):
    api_url = f"https://api-inference.huggingface.co/pipeline/feature-extraction/{HF_AUDIO_MODEL}"
    headers = get_hf_headers()
    with open(audio_path, "rb") as f:
        data = f.read()
        
    response = requests.post(api_url, headers=headers, data=data)
    if response.status_code != 200:
        print(f"Error getting audio embedding: {response.text}")
        return None
    return response.json()

def extract_interval_frames(video_path: Path, interval_sec=10):
    """Extract frames at regular intervals (every 10 seconds)"""
    frames = []
    try:
        cap = cv2.VideoCapture(str(video_path))
        if not cap.isOpened(): return []
        
        fps = cap.get(cv2.CAP_PROP_FPS)
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        duration = total_frames / fps
        
        # Calculate frame indices for every 'interval_sec' seconds
        timestamps = np.arange(0, duration, interval_sec)
        
        for ts in timestamps:
            frame_idx = int(ts * fps)
            cap.set(cv2.CAP_PROP_POS_FRAMES, frame_idx)
            ret, frame = cap.read()
            if ret:
                is_success, buffer = cv2.imencode(".jpg", frame)
                if is_success:
                    frames.append({
                        "timestamp": ts,
                        "data": buffer.tobytes()
                    })
        cap.release()
    except Exception as e:
        print(f"Frame extraction error: {e}")
    return frames

# ... Text extraction functions ...
def extract_text_from_pdf(pdf_path: Path) -> str:
    if not pdfplumber: return ""
    text = ""
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t: text += t + "\n"
    except: pass
    return text

def extract_text_from_docx(docx_path: Path) -> str:
    if not Document: return ""
    try:
        doc = Document(docx_path)
        return "\n".join([p.text for p in doc.paragraphs])
    except: return ""

def extract_text_from_pptx(pptx_path: Path) -> str:
    if not Presentation: return ""
    text = []
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"): text.append(shape.text)
        return "\n".join(text)
    except: return ""

def extract_text_from_txt(txt_path: Path) -> str:
    try: return txt_path.read_text(encoding='utf-8')
    except: return ""

def transcribe_audio(client: OpenAI, audio_path: Path) -> str:
    try:
        with open(audio_path, "rb") as audio_file:
            return client.audio.transcriptions.create(model="whisper-1", file=audio_file).text
    except Exception as e:
        print(f"Transcription error: {e}")
        return ""

def chunk_text(text: str) -> list[str]:
    encoding = tiktoken.get_encoding("cl100k_base")
    tokens = encoding.encode(text)
    chunks = []
    start = 0
    while start < len(tokens):
        end = start + CHUNK_SIZE
        chunk_text = encoding.decode(tokens[start:end])
        chunks.append(chunk_text.strip())
        start = end - CHUNK_OVERLAP
    return [c for c in chunks if c]

def get_openai_embedding(client: OpenAI, text: str) -> list[float]:
    resp = client.embeddings.create(model=OPENAI_TEXT_EMBED_MODEL, input=text)
    return resp.data[0].embedding

def serialize_embedding(embedding: list[float]) -> bytes:
    return struct.pack(f"{len(embedding)}f", *embedding)

def init_database(db_path: Path) -> sqlite3.Connection:
    if db_path.exists(): db_path.unlink()
    conn = sqlite3.connect(str(db_path))
    conn.enable_load_extension(True)
    try:
        conn.load_extension("vss0")
        conn.load_extension("vector0")
    except: pass
    conn.enable_load_extension(False)
    
    # Updated schema with metadata column (JSON or text)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS documents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            file_name TEXT NOT NULL,
            content_type TEXT NOT NULL,
            chunk_text TEXT NOT NULL,
            embedding BLOB NOT NULL,
            embedding_dimensions INTEGER NOT NULL,
            metadata TEXT
        )
    """)
    
    try:
        conn.execute(f"CREATE VIRTUAL TABLE IF NOT EXISTS vss_text USING vss0(embedding({TEXT_EMBED_DIM}))")
        conn.execute(f"CREATE VIRTUAL TABLE IF NOT EXISTS vss_image USING vss0(embedding({IMAGE_EMBED_DIM}))")
        conn.execute(f"CREATE VIRTUAL TABLE IF NOT EXISTS vss_audio USING vss0(embedding({AUDIO_EMBED_DIM}))")
    except: pass
    
    conn.commit()
    return conn

def insert_document(conn, file_name, content_type, chunk_text, embedding, dims, metadata=None):
    if not embedding or not isinstance(embedding, list): return
    
    embedding_bytes = serialize_embedding(embedding)
    cursor = conn.execute(
        "INSERT INTO documents (file_name, content_type, chunk_text, embedding, embedding_dimensions, metadata) VALUES (?, ?, ?, ?, ?, ?)",
        (file_name, content_type, chunk_text, embedding_bytes, dims, metadata)
    )
    doc_id = cursor.lastrowid
    
    try:
        table = "vss_text"
        if content_type == "image": table = "vss_image"
        elif content_type == "audio": table = "vss_audio"
        
        conn.execute(f"INSERT INTO {table} (rowid, embedding) VALUES (?, ?)", (doc_id, embedding_bytes))
    except: pass
    return doc_id

def build_index():
    print(f"Building cloud-first index from {INPUT_DIR}")
    client = get_openai_client()
    conn = init_database(DB_PATH)
    total_chunks = 0
    
    for file_path in INPUT_DIR.iterdir():
        if not file_path.is_file(): continue
        ext = file_path.suffix.lower()
        print(f"\nProcessing {file_path.name}")
        
        # 1. Text Documents
        if ext in DOC_EXTENSIONS:
            text = ""
            if ext == ".pdf": text = extract_text_from_pdf(file_path)
            elif ext == ".docx": text = extract_text_from_docx(file_path)
            elif ext == ".pptx": text = extract_text_from_pptx(file_path)
            else: text = extract_text_from_txt(file_path)
            
            if text.strip():
                chunks = chunk_text(text)
                print(f"  Indexing {len(chunks)} text chunks...")
                for i, chunk in enumerate(chunks):
                    emb = get_openai_embedding(client, chunk)
                    meta = f"chunk_index:{i}"
                    insert_document(conn, file_path.name, "text", chunk, emb, TEXT_EMBED_DIM, meta)
                    total_chunks += 1

        # 2. Images
        elif ext in IMAGE_EXTENSIONS:
            print("  Generating cloud CLIP embedding...")
            emb = get_hf_image_embedding(image_path=file_path)
            if emb:
                insert_document(conn, file_path.name, "image", f"Image: {file_path.name}", emb, IMAGE_EMBED_DIM)
                total_chunks += 1
        
        # 3. Audio / Video
        elif ext in AUDIO_EXTENSIONS or ext in VIDEO_EXTENSIONS:
            # A. Transcribe (Text)
            print("  Transcribing...")
            transcript = transcribe_audio(client, file_path)
            if transcript.strip():
                chunks = chunk_text(transcript)
                for i, chunk in enumerate(chunks):
                    emb = get_openai_embedding(client, chunk)
                    meta = f"chunk_index:{i}"
                    insert_document(conn, file_path.name, "text", f"[Transcript] {chunk}", emb, TEXT_EMBED_DIM, meta)
                    total_chunks += 1
            
            # B. Audio Embedding (CLAP) - for "sound search"
            if ext in AUDIO_EXTENSIONS:
                print("  Generating cloud Audio embedding...")
                emb = get_hf_audio_embedding(file_path)
                if isinstance(emb, list) and len(emb) > 0 and isinstance(emb[0], list):
                    emb = emb[0]
                
                if emb and isinstance(emb, list):
                    insert_document(conn, file_path.name, "audio", f"Audio File: {file_path.name}", emb, AUDIO_EMBED_DIM)
                    total_chunks += 1
            
            # C. Video Visual Frames (CLIP)
            if ext in VIDEO_EXTENSIONS:
                print("  Extracting visual frames (every 10s)...")
                # Every 10 seconds = 6 frames per minute
                frames = extract_interval_frames(file_path, interval_sec=10)
                
                print(f"  Generated {len(frames)} frames. Indexing...")
                for i, frame in enumerate(frames):
                    ts = frame['timestamp']
                    # Embed
                    emb = get_hf_image_embedding(image_data=frame['data'])
                    if emb:
                        # Store metadata: "timestamp:10.5"
                        meta = f"timestamp:{ts:.1f}"
                        # Descriptive text for LLM
                        chunk_text = f"Video Frame: {file_path.name} at {ts:.1f} seconds"
                        
                        insert_document(conn, file_path.name, "image", chunk_text, emb, IMAGE_EMBED_DIM, meta)
                        print(f"    Indexed frame at {ts:.1f}s")
                        total_chunks += 1

    conn.commit()
    conn.close()
    print(f"\nIndexing complete! Total chunks: {total_chunks}")

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--build", action="store_true", help="Build index")
    args = parser.parse_args()
    if args.build: build_index()
    else: parser.print_help()

if __name__ == "__main__":
    main()
